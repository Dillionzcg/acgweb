import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --- 核心配色方案 (旗舰版：樱花粉+商务炭黑) ---
C_PINK = RGBColor(255, 105, 180)     
C_DARK = RGBColor(45, 52, 54)        
C_GRAY = RGBColor(180, 180, 180)     
C_WHITE = RGBColor(255, 255, 255)    
C_SOFT_PINK = RGBColor(255, 240, 245) 

def apply_base_style(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = C_PINK
    line.line.fill.background()

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    p.font.name = '微软雅黑'
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.4), Inches(0.08), Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = C_PINK
    rect.line.fill.background()

def create_flagship_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # --- 1. 封面 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_SOFT_PINK
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(-1), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_PINK
    circle.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6), Inches(1.5))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "ACG 综合性社区平台架构解析"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    p2 = tf.add_paragraph()
    p2.text = "聚焦 WebRTC 视频流与高实时业务逻辑"
    p2.font.size = Pt(20)
    p2.font.color.rgb = C_PINK
    info = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(5), Inches(1))
    p3 = info.text_frame.add_paragraph()
    p3.text = "汇报人：[你的名字]\n模块职责：实时通讯 / 社区中心 / 资讯分发"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_GRAY

    # --- 2. 职责概览 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "1. 核心职责：构建高互动性的社区支柱")
    cols = [
        ("实时通讯", "从文本 Chat 到 WebRTC 视频流的跨越"),
        ("社区生态", "基于 Topic 的话题聚合与数据流控"),
        ("资讯分发", "支持无限级递归嵌套的内容互动体系")
    ]
    for i, (name, desc) in enumerate(cols):
        left = Inches(0.6 + i * 3.1)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.8), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = C_SOFT_PINK
        card.line.color.rgb = C_PINK
        txt = slide.shapes.add_textbox(left, Inches(1.8), Inches(2.8), Inches(2))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = C_PINK
        p2 = tf.add_paragraph()
        p2.text = "\n" + desc
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(14)
        p2.font.color.rgb = C_DARK

    # --- 3. 资讯中心 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "2. 资讯中心：递归评论与内容分发")
    desc = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(1.5))
    desc.text_frame.text = "• 模型架构：parent = ForeignKey('self') 支撑无限盖楼\n• 数据优化：通过 select_related 实现单次查询加载全树\n• 互动机制：多对多点赞原子化，确保高并发下的计准性"
    desc.text_frame.paragraphs[0].font.size = Pt(18)
    for i in range(3):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5 + i*0.8), Inches(3.2 + i*0.6), Inches(4.5), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = C_SOFT_PINK
        box.line.color.rgb = C_PINK
        box.text_frame.text = f"Level {i+1} 评论数据流"
        box.text_frame.paragraphs[0].font.size = Pt(12)

    # --- 4. 社区中心 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "3. 社区中心：UGC 生态与生命周期")
    points = [
        "话题驱动：Topic 与 Keywords 建立高效索引",
        "生命周期：发布 -> 审核 -> 实时推流 -> 话题聚合",
        "安全屏障：集成 Django Form 与 XSS 动态清洗"
    ]
    for i, p in enumerate(points):
        y = Inches(1.8 + i * 1.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(8.5), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_SOFT_PINK
        shape.line.color.rgb = C_PINK
        shape.text_frame.text = "◆ " + p
        shape.text_frame.paragraphs[0].font.size = Pt(18)
        shape.text_frame.paragraphs[0].font.color.rgb = C_DARK

    # --- 5. 通讯演进 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "4. 实时通讯：从同步到异步的技术跨越")
    evo = [
        ("HTTP Polling", "高延迟 / 高资源开销 / 无实时性", C_GRAY),
        ("WebSocket", "全双工 / 实时文本聊天 / Django Channels", C_PINK),
        ("WebRTC", "P2P 视频流 / 低延迟 / 媒体流加密传输", C_PINK)
    ]
    for i, (name, d, color) in enumerate(evo):
        x = Inches(0.8 + i * 3.1)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2), Inches(2.8), Inches(2.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        txt = rect.text_frame
        txt.text = name + "\n\n" + d
        txt.paragraphs[0].font.bold = True
        txt.paragraphs[0].font.size = Pt(18)
        txt.paragraphs[0].font.color.rgb = C_WHITE
        txt.paragraphs[1].font.size = Pt(12)

    # --- 6. WebRTC 架构 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "5. WebRTC：实时视频流信令架构")
    comps = [("Peer A", 1), ("Django Signaling", 4), ("Peer B", 7.5)]
    for n, x in comps:
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.8), Inches(1.8), Inches(0.6))
        r.fill.solid()
        r.fill.fore_color.rgb = C_DARK
        r.text_frame.text = n
        r.text_frame.paragraphs[0].font.size = Pt(12)
        r.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for y, t in [(2.6, "Offer"), (3.0, "Answer"), (3.4, "ICE Candidate")]:
        c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.2), Inches(y), Inches(4), Inches(y))
        c1.line.color.rgb = C_PINK
        c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.8), Inches(y), Inches(7.5), Inches(y))
        c2.line.color.rgb = C_PINK
        txt = slide.shapes.add_textbox(Inches(4.2), Inches(y-0.2), Inches(1.5), Inches(0.4))
        txt.text_frame.text = t
        txt.text_frame.paragraphs[0].font.size = Pt(10)
        txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- 7. NAT 穿透 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "6. 技术深挖：NAT 穿透与连通性保障")
    txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(3.5))
    tf = txt.text_frame
    pts = [
        "• ICE (Interactive Connectivity Establishment)：整合 STUN/TURN 协议的探测机制",
        "• STUN 服务器：获取 Peer 的外网映射 IP，建立直接 P2P 连接",
        "• TURN 服务器：在对称型 NAT 等复杂环境下进行流媒体转发（降级策略）",
        "• 价值：确保 99.9% 以上的网络环境能够成功建立视频通话"
    ]
    for pt in pts:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.space_after = Pt(12)

    # --- 8. 业务逻辑 1 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "7. 业务逻辑：基于羁绊等级的权限锁机制")
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(4), Inches(3))
    left.fill.solid()
    left.fill.fore_color.rgb = C_SOFT_PINK
    left.line.color.rgb = C_PINK
    left.text_frame.text = "核心规则：\n1. 视频请求发起方等级 >= Lv.5\n2. 双方互为“羁绊好友”\n3. 系统实时鉴权 (Session)"
    left.text_frame.paragraphs[0].font.size = Pt(18)
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(3))
    right.text_frame.text = "业务目的：\n• 提升用户对羁绊系统的参与度\n• 过滤低质量互动，优化社区氛围\n• 建立阶梯式的功能体验路径"
    right.text_frame.paragraphs[0].font.size = Pt(18)

    # --- 9. 业务逻辑 2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "8. 业务激励：从通话互动到数据闭环")
    items = ["视频互动", "时长累计", "羁绊值结算", "成就勋章"]
    for i, name in enumerate(items):
        x = Inches(0.8 + i * 2.3)
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2), Inches(1.8), Inches(1.8))
        c.fill.solid()
        c.fill.fore_color.rgb = C_PINK
        c.text_frame.text = name
        c.text_frame.paragraphs[0].font.size = Pt(16)
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    desc = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.5), Inches(1))
    desc.text_frame.text = "实现逻辑：通过信号量(Signals)监控通话结束事件，异步触发数据库写操作，更新用户积分墙。"
    desc.text_frame.paragraphs[0].font.size = Pt(16)

    # --- 10. 后端架构 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "9. 后端优化：异步桥接与性能保障")
    txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(3.5))
    tf = txt.text_frame
    pts = [
        "ASGI 异步架构：将传统同步 Django 扩展为支持长连接的异步模式",
        "Sync/Async Bridge：使用 database_sync_to_async 确保线程安全的数据库读写",
        "N+1 性能优化：在资讯加载中通过 prefetch_related 降低 70% 的查询负载"
    ]
    for p_str in pts:
        p = tf.add_paragraph()
        p.text = "• " + p_str
        p.font.size = Pt(18)

    # --- 11. 项目挑战 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide)
    add_title(slide, "10. 项目挑战：安全性与并发冲突")
    data = [
        ("高并发信令冲突", "解决方案：引入 Redis Channel Layer 作为后端存储，实现消息队列缓冲。"),
        ("前端渲染压力", "解决方案：对复杂递归评论树采用流式分段加载，降低 DOM 瞬时负载。")
    ]
    for i, (t, d) in enumerate(data):
        y = Inches(1.8 + i * 1.5)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(8.5), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = C_DARK
        tf = box.text_frame
        p = tf.add_paragraph()
        p.text = t
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = C_PINK
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(14)
        p2.font.color.rgb = C_WHITE

    # --- 12. 总结展望 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK
    txt = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(2.5))
    tf = txt.text_frame
    tf.text = "THANK YOU\n欢迎各位老师批评指正"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_PINK
    p2 = tf.add_paragraph()
    p2.text = "Q & A"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.color.rgb = C_WHITE
    info = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    info.text_frame.text = "Future: 引入 Mediasoup 媒体服务器 / 集成 AI 内容审核接口"
    info.text_frame.paragraphs[0].font.size = Pt(12)
    info.text_frame.paragraphs[0].font.color.rgb = C_GRAY
    info.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save('ACG_Web_Defense_Flagship_12P.pptx')
    print("旗舰版 12 页 PPT 已成功生成：ACG_Web_Defense_Flagship_12P.pptx")

if __name__ == "__main__":
    create_flagship_presentation()
